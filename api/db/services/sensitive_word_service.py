#
#  Copyright 2026 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
"""Sensitive-word management + a fast upload-time content scanner.

The scanner is deliberately cheap: it does a *fast text extraction* only (no
OCR, no layout analysis) so it barely adds to upload latency, then does a plain
substring match. Any file type that would require OCR (scanned PDF, images) is
passed through untouched. Everything here is fail-open: any extraction/match
error is swallowed and treated as "no hit" so the scan can never break uploads.
"""
import logging
import time
from io import BytesIO
from pathlib import Path

from api.db.db_models import DB, SensitiveWord
from api.db.services.common_service import CommonService

# Feature flag stored in the shared ``system_settings`` table.
SENSITIVE_FILTER_ENABLED_KEY = "kb.sensitive_filter.enabled"

# Machine-readable prefix embedded in the per-file upload error so the frontend
# can recognize a sensitive-word rejection and localize it. Followed by the
# comma-joined hit words, e.g. ``__SENSITIVE_WORDS__:foo,bar``.
SENSITIVE_HIT_PREFIX = "__SENSITIVE_WORDS__:"

# Bound the work per file so a huge document can never stall the upload:
#   - only look at the first N pages of a PDF
#   - only match against the first N characters of extracted text
_PDF_MAX_PAGES = 100
_MAX_SCAN_CHARS = 2_000_000

# Plain-text-ish suffixes we can scan by simply decoding the bytes.
_PLAIN_TEXT_SUFFIXES = {
    "txt", "md", "markdown", "csv", "tsv", "json", "log",
    "html", "htm", "xml", "yaml", "yml", "ini", "conf",
}


class SensitiveWordService(CommonService):
    model = SensitiveWord

    @classmethod
    @DB.connection_context()
    def list_words(cls):
        return list(cls.model.select().order_by(cls.model.create_time.desc()))

    @classmethod
    @DB.connection_context()
    def word_exists(cls, word: str) -> bool:
        return cls.model.select().where(cls.model.word == word).count() > 0

    @classmethod
    @DB.connection_context()
    def all_words(cls) -> list[str]:
        return [r.word for r in cls.model.select(cls.model.word)]


# --------------------------------------------------------------------------- #
# In-process cache so the upload path does not hit the DB once per file.
# --------------------------------------------------------------------------- #
_CACHE_TTL = 30  # seconds
_cache: dict = {"ts": 0.0, "enabled": False, "words": []}


def _load_config() -> tuple[bool, list[str]]:
    now = time.time()
    if now - _cache["ts"] < _CACHE_TTL:
        return _cache["enabled"], _cache["words"]

    enabled, words = False, []
    try:
        from api.db.services.system_settings_service import SystemSettingsService

        rows = SystemSettingsService.get_by_name(SENSITIVE_FILTER_ENABLED_KEY)
        rows = list(rows)
        if rows:
            enabled = str(rows[0].value).strip().lower() == "true"
    except Exception:
        logging.exception("Failed to read sensitive-filter enabled flag")

    if enabled:
        try:
            words = [w for w in SensitiveWordService.all_words() if w]
        except Exception:
            logging.exception("Failed to load sensitive words")
            words = []

    _cache.update(ts=now, enabled=enabled, words=words)
    return enabled, words


def invalidate_cache() -> None:
    """Drop the cache so the next scan reloads config (call after CRUD)."""
    _cache["ts"] = 0.0


# --------------------------------------------------------------------------- #
# Fast, best-effort text extraction (no OCR).
# --------------------------------------------------------------------------- #
def _extract_text_fast(filename: str, blob: bytes) -> str:
    suffix = Path(filename or "").suffix.lstrip(".").lower()
    try:
        if suffix in _PLAIN_TEXT_SUFFIXES:
            from rag.nlp import find_codec

            return blob.decode(find_codec(blob), errors="ignore")

        if suffix == "docx":
            from docx import Document

            doc = Document(BytesIO(blob))
            return "\n".join(p.text for p in doc.paragraphs)

        if suffix == "xlsx":
            import openpyxl

            wb = openpyxl.load_workbook(BytesIO(blob), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            parts.append(str(cell))
                    if sum(len(p) for p in parts) > _MAX_SCAN_CHARS:
                        break
            wb.close()
            return " ".join(parts)

        if suffix == "pptx":
            from pptx import Presentation

            prs = Presentation(BytesIO(blob))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)

        if suffix == "pdf":
            # Text-based PDFs only; scanned PDFs yield ~empty text and are
            # passed through (no OCR at upload time).
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(blob))
            parts = []
            for page in reader.pages[:_PDF_MAX_PAGES]:
                parts.append(page.extract_text() or "")
                if sum(len(p) for p in parts) > _MAX_SCAN_CHARS:
                    break
            return "\n".join(parts)
    except Exception:
        # Fail open: any extraction problem is treated as "nothing to scan".
        logging.debug("Sensitive-word fast extraction failed for %s", filename, exc_info=True)
        return ""

    # Images and every other type require OCR to read -> pass through.
    return ""


def scan_blob(filename: str, blob: bytes) -> list[str]:
    """Return the list of sensitive words found in ``blob``.

    An empty list means "allow the upload" (feature off, no words configured,
    nothing extractable, or genuinely clean). Never raises.
    """
    try:
        enabled, words = _load_config()
        if not enabled or not words:
            return []

        text = _extract_text_fast(filename, blob)
        if not text:
            return []
        text = text[:_MAX_SCAN_CHARS]

        # Case-insensitive for ASCII (so English words match regardless of
        # case); Chinese/full-width text is matched as-is.
        haystack = text.lower()
        hits = []
        for w in words:
            if w.lower() in haystack:
                hits.append(w)
        return hits
    except Exception:
        logging.exception("Sensitive-word scan failed for %s", filename)
        return []
